#!/usr/bin/env python3
"""
Script to generate a pitch deck PowerPoint presentation from structured data.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck(data, output_file="pitch_deck.pptx"):
    """
    Create a pitch deck from structured data.

    Args:
        data: Dictionary containing pitch deck information
        output_file: Path to save the PowerPoint file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define Earthen Brand color scheme
    PRIMARY_COLOR = RGBColor(45, 81, 55)  # Forest Green #2D5137
    SECONDARY_COLOR = RGBColor(191, 109, 73)  # Terracotta #BF6D49
    ACCENT_COLOR = RGBColor(242, 199, 68)  # Sunlight #F2C744
    TEXT_COLOR = RGBColor(51, 51, 51)  # Charcoal #333333
    SAND_COLOR = RGBColor(245, 241, 234)  # Sand #F5F1EA

    def add_title_slide(title, subtitle=""):
        """Add a title slide with Earthen Brand styling"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add sand background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SAND_COLOR

        # Add decorative accent bar
        accent_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(3.2), Inches(10), Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT_COLOR
        accent_bar.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(26)
            subtitle_para.font.color.rgb = TEXT_COLOR
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_items):
        """Add a content slide with title and bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add sand background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SAND_COLOR

        # Add header background shape
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(1.3)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Add decorative accent element
        accent_circle = slide.shapes.add_shape(
            3,  # Oval
            Inches(8.8), Inches(0.2), Inches(0.9), Inches(0.9)
        )
        accent_circle.fill.solid()
        accent_circle.fill.fore_color.rgb = ACCENT_COLOR
        accent_circle.line.fill.background()

        # Add content with rounded background
        content_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.8), Inches(9), Inches(5.2)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_bg.line.color.rgb = RGBColor(200, 200, 200)
        content_bg.line.width = Pt(1)

        # Add content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(14)
            p.level = 0

        return slide

    def add_two_column_slide(title, left_title, left_content, right_title, right_content):
        """Add a two-column content slide with visual separation"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add sand background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SAND_COLOR

        # Add header background shape
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(1.3)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Left column background
        left_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.2)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        left_bg.line.color.rgb = SECONDARY_COLOR
        left_bg.line.width = Pt(2)

        # Right column background
        right_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.2)
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        right_bg.line.color.rgb = PRIMARY_COLOR
        right_bg.line.width = Pt(2)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(3.7), Inches(4.6))
        left_frame = left_box.text_frame
        left_frame.text = left_title
        left_para = left_frame.paragraphs[0]
        left_para.font.size = Pt(22)
        left_para.font.bold = True
        left_para.font.color.rgb = SECONDARY_COLOR

        for item in left_content:
            left_frame.add_paragraph()
            p = left_frame.paragraphs[-1]
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.1), Inches(3.7), Inches(4.6))
        right_frame = right_box.text_frame
        right_frame.text = right_title
        right_para = right_frame.paragraphs[0]
        right_para.font.size = Pt(22)
        right_para.font.bold = True
        right_para.font.color.rgb = PRIMARY_COLOR

        for item in right_content:
            right_frame.add_paragraph()
            p = right_frame.paragraphs[-1]
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

        return slide

    # Slide 1: Title slide
    company_name = data.get("company_name", "Company Name")
    tagline = data.get("tagline", "")
    add_title_slide(company_name, tagline)

    # Slide 2: Problem
    if "problem" in data:
        problem_items = data["problem"] if isinstance(data["problem"], list) else [data["problem"]]
        add_content_slide("The Problem", problem_items)

    # Slide 3: Solution
    if "solution" in data:
        solution_items = data["solution"] if isinstance(data["solution"], list) else [data["solution"]]
        add_content_slide("Our Solution", solution_items)

    # Slide 4: Market Opportunity
    if "market" in data:
        market_items = data["market"] if isinstance(data["market"], list) else [data["market"]]
        add_content_slide("Market Opportunity", market_items)

    # Slide 5: Product
    if "product" in data:
        product_items = data["product"] if isinstance(data["product"], list) else [data["product"]]
        add_content_slide("Product", product_items)

    # Slide 6: Traction
    if "traction" in data:
        traction_items = data["traction"] if isinstance(data["traction"], list) else [data["traction"]]
        add_content_slide("Traction", traction_items)

    # Slide 7: Business Model
    if "business_model" in data:
        bm_items = data["business_model"] if isinstance(data["business_model"], list) else [data["business_model"]]
        add_content_slide("Business Model", bm_items)

    # Slide 8: Competition
    if "competition" in data:
        comp_data = data["competition"]
        if isinstance(comp_data, dict) and "our_advantages" in comp_data and "competitors" in comp_data:
            add_two_column_slide(
                "Competitive Landscape",
                "Our Advantages",
                comp_data.get("our_advantages", []),
                "Competition",
                comp_data.get("competitors", [])
            )
        else:
            comp_items = comp_data if isinstance(comp_data, list) else [comp_data]
            add_content_slide("Competitive Landscape", comp_items)

    # Slide 9: Team
    if "team" in data:
        team_items = data["team"] if isinstance(data["team"], list) else [data["team"]]
        add_content_slide("Team", team_items)

    # Slide 10: Financials / Ask
    if "financials" in data:
        fin_items = data["financials"] if isinstance(data["financials"], list) else [data["financials"]]
        add_content_slide("Financials & Ask", fin_items)

    # Save presentation
    prs.save(output_file)
    print(f"✅ Pitch deck created: {output_file}")
    return output_file


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 create_pitch_deck.py <data.json> [output.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "pitch_deck.pptx"

    with open(input_file, 'r') as f:
        data = json.load(f)

    create_pitch_deck(data, output_file)
